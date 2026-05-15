# pattern: Imperative Shell
"""Convert DOCX and PPTX files to page dicts using python-docx / python-pptx (torch-free).

Page dicts match the shape produced by parse_pdf.py so all downstream
chunking (chunk_templates.py) works without modification.

Bboxes are synthetic (full letter-page size) — DOCX/PPTX have no per-line
pixel coordinates. The SourcePanel will gracefully show 'Could not load PDF'
for these document types; that is acceptable for this version.
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

log = logging.getLogger(__name__)

_SYNTHETIC_PAGE_BBOX = [0.0, 0.0, 612.0, 792.0]
_CHARS_PER_PAGE = 3000


def load_document_with_docling(file_path: str) -> list[dict[str, Any]]:
    """Convert a DOCX or PPTX file to a list of page dicts.

    Each dict has: page_number, text, tables, confidence, bbox, text_lines, parser.
    Large documents are split into ~3000-char synthetic pages for downstream chunking.
    """
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".docx":
        full_text = _extract_docx(file_path)
    elif suffix == ".pptx":
        full_text = _extract_pptx(file_path)
    else:
        raise ValueError(f"Unsupported format for docling_loader: {suffix}")

    if not full_text.strip():
        log.warning("Produced empty text for '%s'", file_path)
        return [_make_page(1, "(empty document)")]

    return _split_into_pages(full_text)


def _extract_docx(file_path: str) -> str:
    import docx  # python-docx

    doc = docx.Document(file_path)
    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)
    return "\n\n".join(parts)


def _extract_pptx(file_path: str) -> str:
    from pptx import Presentation  # python-pptx

    prs = Presentation(file_path)
    parts: list[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)
        if slide_parts:
            parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_parts))
    return "\n\n".join(parts)


def _split_into_pages(text: str) -> list[dict[str, Any]]:
    """Break text into synthetic pages of ~_CHARS_PER_PAGE characters, splitting at paragraph boundaries."""
    pages = []
    paragraphs = text.split("\n\n")
    current_parts: list[str] = []
    current_len = 0
    page_num = 1

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        if current_len + len(para) > _CHARS_PER_PAGE and current_parts:
            pages.append(_make_page(page_num, "\n\n".join(current_parts)))
            current_parts = []
            current_len = 0
            page_num += 1
        current_parts.append(para)
        current_len += len(para)

    if current_parts:
        pages.append(_make_page(page_num, "\n\n".join(current_parts)))

    return pages


def _make_page(page_number: int, text: str) -> dict[str, Any]:
    return {
        "page_number": page_number,
        "text": text,
        "tables": [],
        "confidence": 100.0,
        "bbox": list(_SYNTHETIC_PAGE_BBOX),
        "text_lines": [],
        "parser": "python-docx/pptx",
    }
